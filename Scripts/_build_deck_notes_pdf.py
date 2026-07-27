#!/usr/bin/env python3
"""Build a read-aloud notes PDF from a study companion deck.

One page per slide: the slide image on top, the slide's speaker-notes script
below. Nothing else — the primary-text background and Q&A material live in the
Presenter's Companion, not here.

Why this exists rather than PowerPoint's own notes-pages export: PowerPoint's
``ExportAsFixedFormat`` (the only API that accepts an ``OutputType`` of
``ppPrintOutputNotesPages``) is not reachable in this environment — every arity
raises "The Python instance can not be converted to a COM object", which is why
``_pptx_to_pdf.py`` silently falls back to ``SaveAs``, and ``SaveAs`` has no
``OutputType`` parameter. Native notes pages also *clip* text that overflows the
notes placeholder; these scripts run to ~2,500 characters, so clipping would
silently drop the end of a presenter's script. Composing the pages here keeps
every word and paginates when a script is long.

Speaker notes are read from the .pptx itself, so this always reflects what is
actually in the deck. Slide images come from the deck PDF, which is generated
on demand if it is missing or stale.

Examples (from repo root):

  python Scripts/_build_deck_notes_pdf.py --study The-Ontology-of-Coexistence
  python Scripts/_build_deck_notes_pdf.py Studies/<Slug>/Deck.pptx
  python Scripts/_build_deck_notes_pdf.py Studies/<Slug>/Deck.pptx -o /tmp/notes.pdf
"""
from __future__ import annotations

import argparse
import re
import sys
from pathlib import Path

import fitz  # PyMuPDF
from pptx import Presentation

from _common import STUDIES
from _pptx_to_pdf import convert_pptx_to_pdf, resolve_pptx

# --- page geometry (A4 portrait, points) ------------------------------------
PAGE_W, PAGE_H = 595.28, 841.89
MARGIN = 46.0
CONTENT_W = PAGE_W - 2 * MARGIN
HEADER_Y = 56.0
IMAGE_TOP = 78.0
BODY_GAP = 24.0
BOTTOM_LIMIT = PAGE_H - 52.0

# --- type ------------------------------------------------------------------
BODY_SIZE = 9.6
BODY_LEAD = 13.2
PARA_GAP = 6.4
HEAD_SIZE = 13.0
KICKER_SIZE = 8.2
FOOT_SIZE = 8.0

# --- palette (deck design tokens) ------------------------------------------
NAVY = (0x1E / 255, 0x24 / 255, 0x47 / 255)
BODY_INK = (0x2B / 255, 0x30 / 255, 0x40 / 255)
GREY = (0x5A / 255, 0x63 / 255, 0x77 / 255)
GOLD = (0x9A / 255, 0x64 / 255, 0x14 / 255)
RULE = (0xC9 / 255, 0xD4 / 255, 0xEF / 255)

FONT_FILES = {
    "body": r"C:\Windows\Fonts\calibri.ttf",
    "bold": r"C:\Windows\Fonts\calibrib.ttf",
    "italic": r"C:\Windows\Fonts\calibrii.ttf",
    "head": r"C:\Windows\Fonts\cambriab.ttf",
}
FALLBACK_FONTS = {"body": "helv", "bold": "hebo", "italic": "heit", "head": "hebo"}


def load_fonts() -> dict[str, fitz.Font]:
    fonts: dict[str, fitz.Font] = {}
    for key, path in FONT_FILES.items():
        if Path(path).is_file():
            fonts[key] = fitz.Font(fontfile=path)
        else:
            fonts[key] = fitz.Font(FALLBACK_FONTS[key])
    return fonts


# --- markdown-ish inline emphasis ------------------------------------------
TOKEN_RE = re.compile(r"(\*\*[^*]+\*\*|(?<!\*)\*[^*]+\*(?!\*))")


def tokenize(paragraph: str) -> list[tuple[str, str, str]]:
    """Split a paragraph into (kind, text, style) tokens.

    kind is 'ws' or 'word'; style is 'body', 'bold' or 'italic'. Whitespace is
    kept as its own token so that punctuation glued to an emphasised word — as
    in "*worked out*." — is never pushed onto the next line with a space.
    """
    tokens: list[tuple[str, str, str]] = []
    for chunk in TOKEN_RE.split(paragraph):
        if not chunk:
            continue
        if chunk.startswith("**") and chunk.endswith("**"):
            style, text = "bold", chunk[2:-2]
        elif chunk.startswith("*") and chunk.endswith("*"):
            style, text = "italic", chunk[1:-1]
        else:
            style, text = "body", chunk
        for piece in re.split(r"(\s+)", text):
            if not piece:
                continue
            tokens.append(("ws" if piece.isspace() else "word", piece, style))
    return tokens


def wrap_paragraph(
    paragraph: str, fonts: dict[str, fitz.Font], width: float
) -> list[list[tuple[str, str]]]:
    """Wrap one paragraph into lines of (text, style) segments."""
    space_w = fonts["body"].text_length(" ", fontsize=BODY_SIZE)
    lines: list[list[tuple[str, str]]] = []
    line: list[tuple[str, str]] = []
    line_w = 0.0
    pending_space = False

    for kind, text, style in tokenize(paragraph):
        if kind == "ws":
            if line:
                pending_space = True
            continue
        word_w = fonts[style].text_length(text, fontsize=BODY_SIZE)
        extra = space_w if pending_space else 0.0
        if line and line_w + extra + word_w > width:
            lines.append(line)
            line, line_w, pending_space = [], 0.0, False
            extra = 0.0
        if pending_space:
            line.append((" ", "body"))
            line_w += space_w
            pending_space = False
        line.append((text, style))
        line_w += word_w
    if line:
        lines.append(line)
    return lines


def layout_note(
    note: str, fonts: dict[str, fitz.Font], width: float
) -> list[list[list[tuple[str, str]]]]:
    """Wrap a whole note into paragraphs of lines.

    Split on any run of newlines: python-pptx returns one newline per notes
    paragraph, while the companion markdown separates them with a blank line.
    Both must yield the same paragraphs here.
    """
    paragraphs = [p.strip() for p in re.split(r"\n+", note) if p.strip()]
    return [wrap_paragraph(re.sub(r"\s+", " ", p), fonts, width) for p in paragraphs]


# --- deck reading ----------------------------------------------------------
def slide_title(slide) -> str:
    """Best-effort slide title: the text run with the largest font size."""
    best, best_size = "", -1.0
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                if not run.text.strip():
                    continue
                size = run.font.size.pt if run.font.size else 0.0
                if size > best_size:
                    best, best_size = run.text.strip(), size
    return re.sub(r"\s+", " ", best)


def read_deck(pptx: Path) -> list[tuple[str, str]]:
    """Return [(title, notes)] in presentation order."""
    prs = Presentation(str(pptx))
    out: list[tuple[str, str]] = []
    for slide in prs.slides:
        note = ""
        if slide.has_notes_slide:
            note = slide.notes_slide.notes_text_frame.text.strip()
        out.append((slide_title(slide), note))
    return out


def ensure_deck_pdf(pptx: Path, deck_pdf: Path | None) -> Path:
    pdf = deck_pdf or pptx.with_suffix(".pdf")
    if pdf.is_file() and pdf.stat().st_mtime >= pptx.stat().st_mtime:
        return pdf
    convert_pptx_to_pdf(pptx, pdf)
    return pdf


# --- drawing ---------------------------------------------------------------
def draw_header(page, fonts, slide_no: int, total: int, title: str, cont: bool) -> None:
    tw = fitz.TextWriter(page.rect)
    kicker = "SLIDE %d OF %d%s" % (slide_no, total, "  ·  CONTINUED" if cont else "")
    tw.append(fitz.Point(MARGIN, HEADER_Y - 14), kicker, font=fonts["body"],
              fontsize=KICKER_SIZE)
    tw.write_text(page, color=GOLD)

    tw = fitz.TextWriter(page.rect)
    shown = title
    while (fonts["head"].text_length(shown, fontsize=HEAD_SIZE) > CONTENT_W
           and len(shown) > 8):
        shown = shown[:-2]
    if shown != title:
        shown = shown.rstrip() + "…"
    tw.append(fitz.Point(MARGIN, HEADER_Y + 2), shown, font=fonts["head"],
              fontsize=HEAD_SIZE)
    tw.write_text(page, color=NAVY)


def draw_footer(page, fonts, page_no: int, deck_name: str) -> None:
    tw = fitz.TextWriter(page.rect)
    left = "%s — read-aloud notes" % deck_name
    tw.append(fitz.Point(MARGIN, PAGE_H - 30), left, font=fonts["body"],
              fontsize=FOOT_SIZE)
    num = str(page_no)
    num_w = fonts["body"].text_length(num, fontsize=FOOT_SIZE)
    tw.append(fitz.Point(PAGE_W - MARGIN - num_w, PAGE_H - 30), num,
              font=fonts["body"], fontsize=FOOT_SIZE)
    tw.write_text(page, color=GREY)


def draw_lines(page, fonts, paragraphs, start_y: float) -> tuple[int, int, float]:
    """Draw as many paragraphs/lines as fit. Returns (para_idx, line_idx, y)."""
    tw = fitz.TextWriter(page.rect)
    y = start_y
    drawn_any = False
    for pi, para in enumerate(paragraphs):
        for li, line in enumerate(para):
            if y > BOTTOM_LIMIT:
                if drawn_any:
                    tw.write_text(page, color=BODY_INK)
                return pi, li, y
            x = MARGIN
            for text, style in line:
                tw.append(fitz.Point(x, y), text, font=fonts[style],
                          fontsize=BODY_SIZE)
                x += fonts[style].text_length(text, fontsize=BODY_SIZE)
            drawn_any = True
            y += BODY_LEAD
        y += PARA_GAP
    if drawn_any:
        tw.write_text(page, color=BODY_INK)
    return len(paragraphs), 0, y


def build(pptx: Path, deck_pdf: Path, out_pdf: Path) -> tuple[int, int]:
    fonts = load_fonts()
    slides = read_deck(pptx)
    total = len(slides)
    src = fitz.open(str(deck_pdf))
    if len(src) != total:
        raise SystemExit(
            "Deck PDF has %d pages but the .pptx has %d slides — regenerate the "
            "deck PDF first." % (len(src), total)
        )

    doc = fitz.open()
    deck_name = pptx.stem.replace("-", " ")
    page_no = 0

    for idx, (title, note) in enumerate(slides):
        heading = title or "Slide %d" % (idx + 1)
        paragraphs = layout_note(note, fonts, CONTENT_W) if note else []

        page = doc.new_page(width=PAGE_W, height=PAGE_H)
        page_no += 1
        draw_header(page, fonts, idx + 1, total, heading, cont=False)

        # Slide, scaled to the deck's own aspect ratio. Placed as vector via
        # show_pdf_page rather than a rasterised pixmap: sharp at any zoom and
        # roughly a sixth of the file size.
        src_page = src[idx]
        ratio = src_page.rect.height / src_page.rect.width
        img_h = CONTENT_W * ratio
        rect = fitz.Rect(MARGIN, IMAGE_TOP, MARGIN + CONTENT_W, IMAGE_TOP + img_h)
        page.show_pdf_page(rect, src, idx)
        page.draw_rect(rect, color=RULE, width=0.5)

        y = rect.y1 + BODY_GAP
        if not paragraphs:
            tw = fitz.TextWriter(page.rect)
            tw.append(fitz.Point(MARGIN, y), "(no read-aloud script for this slide)",
                      font=fonts["italic"], fontsize=BODY_SIZE)
            tw.write_text(page, color=GREY)
            draw_footer(page, fonts, page_no, deck_name)
            continue

        pi, li, _ = draw_lines(page, fonts, paragraphs, y)
        draw_footer(page, fonts, page_no, deck_name)

        # Continuation pages: no slide image, full column, until the script ends.
        while pi < len(paragraphs):
            remaining = [paragraphs[pi][li:]] + list(paragraphs[pi + 1:])
            page = doc.new_page(width=PAGE_W, height=PAGE_H)
            page_no += 1
            draw_header(page, fonts, idx + 1, total, heading, cont=True)
            npi, nli, _ = draw_lines(page, fonts, remaining, IMAGE_TOP + 6)
            draw_footer(page, fonts, page_no, deck_name)
            if npi == 0 and nli == li:
                raise SystemExit(
                    "No progress paginating slide %d — check the layout constants."
                    % (idx + 1)
                )
            pi, li = (pi + npi, nli) if npi else (pi, nli)

    out_pdf.parent.mkdir(parents=True, exist_ok=True)
    # Subsetting the four embedded faces and collecting garbage takes this from
    # ~3.5 MB to well under 1 MB; the deck pages are already vector.
    try:
        doc.subset_fonts()
    except Exception:  # pragma: no cover - older PyMuPDF without subset_fonts
        pass
    doc.save(str(out_pdf), deflate=True, deflate_fonts=True, deflate_images=True,
             garbage=4, clean=True)
    doc.close()
    src.close()
    return total, page_no


def main(argv: list[str] | None = None) -> int:
    parser = argparse.ArgumentParser(
        description="Build a read-aloud notes PDF (slide + speaker script per page) "
        "from a companion deck.",
    )
    parser.add_argument("pptx", nargs="?", type=Path, help="Path to the .pptx file")
    parser.add_argument("--study", help="Study slug under Studies/")
    parser.add_argument("--deck", help="Deck filename under Studies/<Slug>/")
    parser.add_argument(
        "--deck-pdf",
        type=Path,
        help="Slide-image source PDF (default: the deck's own .pdf, generated if stale)",
    )
    parser.add_argument(
        "--output", "-o", type=Path,
        help="Output .pdf path (default: <deck stem>-notes.pdf beside the deck)",
    )
    args = parser.parse_args(argv)

    pptx = resolve_pptx(args.pptx, args.study, args.deck)
    out = args.output.expanduser().resolve() if args.output else (
        pptx.with_name(pptx.stem + "-notes.pdf")
    )
    if out.suffix.lower() != ".pdf":
        raise SystemExit("Output must end with .pdf: %s" % out)

    deck_pdf = ensure_deck_pdf(pptx, args.deck_pdf)
    slides, pages = build(pptx, deck_pdf, out)
    print("Wrote %s (%d slides, %d pages)" % (out, slides, pages))
    return 0


if __name__ == "__main__":
    raise SystemExit(main())

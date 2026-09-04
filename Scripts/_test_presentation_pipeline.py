#!/usr/bin/env python3
"""Regression tests for the presentation manifest and artifact verifier."""
from __future__ import annotations

import tempfile
import unittest
from pathlib import Path

import pymupdf as fitz
from pptx import Presentation
from pptx.util import Inches

from _presentation_pipeline import DeckSpec, load_manifest, manifest_errors
from _pptx_to_pdf import resolve_output
from _verify_presentation_reproducible import compare_artifact
from _verify_presentations import token_recall, verify_slides_pdf


class PresentationManifestTests(unittest.TestCase):
    def test_repository_manifest_covers_every_deck(self) -> None:
        manifest = load_manifest()
        self.assertEqual(manifest_errors(manifest), [])
        self.assertEqual(len(manifest.decks), 7)
        self.assertEqual(manifest.profile().status, "accepted-production")

    def test_same_basename_decks_do_not_collide_with_study_pdf(self) -> None:
        manifest = load_manifest()
        for deck_id in ("undivided-society", "not-just-material"):
            spec = manifest.deck(deck_id)
            canonical = spec.source.parent / f"{spec.source.parent.name}.pdf"
            self.assertNotEqual(spec.slides_pdf, canonical)
            self.assertEqual(resolve_output(spec.source, None), spec.slides_pdf)

    def test_every_slide_has_speaker_notes(self) -> None:
        for spec in load_manifest().decks:
            prs = Presentation(str(spec.source))
            missing = [
                index
                for index, slide in enumerate(prs.slides, 1)
                if not slide.has_notes_slide
                or not slide.notes_slide.notes_text_frame.text.strip()
            ]
            self.assertEqual(missing, [], f"{spec.id} has slides without notes")


class PresentationVerifierTests(unittest.TestCase):
    def test_token_recall_counts_repeated_words(self) -> None:
        self.assertEqual(token_recall("one one two", "one two"), 2 / 3)
        self.assertEqual(token_recall("", "anything"), 1.0)

    def test_wrong_page_geometry_and_missing_text_fail(self) -> None:
        with tempfile.TemporaryDirectory() as temp:
            root = Path(temp)
            pptx = root / "sample.pptx"
            pdf = root / "sample.pdf"
            prs = Presentation()
            prs.slide_width = Inches(10)
            prs.slide_height = Inches(5.625)
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
            box.text = "Required presentation words"
            prs.save(pptx)

            doc = fitz.open()
            page = doc.new_page(width=595.28, height=841.89)
            page.insert_text((72, 72), "Different content")
            doc.save(pdf)
            doc.close()

            spec = DeckSpec("sample", pptx, pdf, root / "notes.pdf", ("Calibri",))
            errors = verify_slides_pdf(spec, Presentation(str(pptx)), pdf)
            self.assertTrue(any("aspect ratio" in error for error in errors))
            self.assertTrue(any("preserves only" in error for error in errors))
            self.assertTrue(any("required font family" in error for error in errors))

    def test_reproducibility_ignores_metadata_but_not_rendered_content(self) -> None:
        with tempfile.TemporaryDirectory() as temp:
            root = Path(temp)
            left = root / "left.pdf"
            right = root / "right.pdf"
            changed = root / "changed.pdf"

            for path, title, text in (
                (left, "First metadata", "Stable content"),
                (right, "Second metadata", "Stable content"),
                (changed, "Second metadata", "Changed content"),
            ):
                doc = fitz.open()
                page = doc.new_page(width=720, height=405)
                page.insert_text((72, 72), text)
                doc.set_metadata({"title": title})
                doc.save(path, no_new_id=True)
                doc.close()

            self.assertNotEqual(left.read_bytes(), right.read_bytes())
            self.assertEqual(compare_artifact("sample", left, right), [])
            self.assertTrue(compare_artifact("sample", left, changed))


if __name__ == "__main__":
    unittest.main()

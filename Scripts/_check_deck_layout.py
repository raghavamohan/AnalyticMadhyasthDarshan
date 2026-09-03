#!/usr/bin/env python3
"""Check a companion deck for layout faults that only show up once it is rendered.

Usage (from repo root):
    python Scripts/_check_deck_layout.py Studies/<Slug>/<Deck>.pptx
    python Scripts/_check_deck_layout.py --study The-Epistemology-of-Coexistence --deck <Deck>.pptx
    python Scripts/_check_deck_layout.py --all

PowerPoint does not clip text that outgrows its box -- it spills, and whether
that spill is a defect depends entirely on what sits next to it. So the fatal
checks here are about where the *rendered* text actually lands, not about
whether it stayed inside its declared box.

Checks:
  1. Rendered text colliding with another shape's rendered text. This is how a
     too-long title shows up: the title box is one line tall and centre-
     anchored, so a wrapped title grows in both directions and its first line
     rides up through the eyebrow above it.
  2. Rendered text spilling off the slide canvas.
  3. "N / M" slide-number footers: N must match the slide's position and M the
     deck's slide count. Adding or removing a slide invalidates every one.
  4. Hard-coded "slide N" cross-references in shape text and speaker notes.
     Out-of-range targets fail; in-range ones are listed, because a reorder
     silently repoints them at the wrong slide.
  5. Reported as notes, not failures: text needing more lines than its declared
     box holds (usually harmless, occasionally the first sign of a fault), and
     single-line boxes at or above --warn-fill of their width.

Widths are measured with PIL against the real font file, which agrees with
PowerPoint to within about a percent -- close enough to catch a wrap, not close
enough to adjudicate a box at 99% fill. That margin is what the note tier is
for: a title measuring 100.2% here still rendered on one line in PowerPoint.

Exit code 1 when any check fails.
"""
from __future__ import annotations

import argparse
import math
import re
from dataclasses import dataclass
from pathlib import Path

from PIL import ImageFont
from pptx import Presentation
from pptx.util import Emu

from _common import STUDIES, configure_utf8_stdio
from _pptx_to_pdf import resolve_pptx

# Windows ships these under C:\Windows\Fonts; the decks use Cambria and Calibri
# throughout, with Courier New for the occasional code-ish label.
FONT_DIR = Path(r"C:\Windows\Fonts")
FONT_FILES: dict[tuple[str, bool, bool], str] = {
    ("cambria", False, False): "cambria.ttc",
    ("cambria", True, False): "cambriab.ttf",
    ("cambria", False, True): "cambriai.ttf",
    ("cambria", True, True): "cambriaz.ttf",
    ("calibri", False, False): "calibri.ttf",
    ("calibri", True, False): "calibrib.ttf",
    ("calibri", False, True): "calibrii.ttf",
    ("calibri", True, True): "calibriz.ttf",
    ("courier new", False, False): "cour.ttf",
    ("courier new", True, False): "courbd.ttf",
    ("courier new", False, True): "couri.ttf",
    ("courier new", True, True): "courbi.ttf",
    ("times new roman", False, False): "times.ttf",
    ("times new roman", True, False): "timesbd.ttf",
    ("arial", False, False): "arial.ttf",
    ("arial", True, False): "arialbd.ttf",
}

DEFAULT_FONT = "Calibri"
DEFAULT_SIZE_PT = 14.0
# PowerPoint's single-spaced line box for these faces is about 1.2 em.
LINE_HEIGHT = 1.2
# Measure at 10x and divide, so integer font sizes don't round away the detail.
SCALE = 10
# PIL and PowerPoint disagree by well under a percent on these faces. Inside that
# band, believe PowerPoint and keep the line unbroken -- a title measuring 100.2%
# here does render on one line. The --warn-fill note tier covers the same band, so
# nothing in it goes unreported; it just isn't called a failure.
FIT_TOLERANCE = 1.01
# Line placement is modelled, not rendered, and the error accumulates down a
# block, so a hairline overlap is not evidence of anything. Past this much the
# overlap is larger than the model's error and every instance checked by hand was
# a real collision; below it, say so and let a human look. Do not raise this to
# silence a finding -- render the slide instead.
COLLISION_FATAL_IN = 0.08

SLIDE_REF = re.compile(r"\bslides?\s+(\d+)\b", re.IGNORECASE)
PAGE_NUMBER = re.compile(r"(\d+)\s*/\s*(\d+)\s*$")

_font_cache: dict[tuple[str, int], ImageFont.FreeTypeFont] = {}
_missing_fonts: set[str] = set()
_rotated: set[str] = set()


@dataclass
class Finding:
    slide: int
    kind: str
    detail: str
    fatal: bool = True


def load_font(name: str, bold: bool, italic: bool, size_pt: float):
    key = (name or DEFAULT_FONT).strip().lower()
    filename = (
        FONT_FILES.get((key, bool(bold), bool(italic)))
        or FONT_FILES.get((key, bool(bold), False))
        or FONT_FILES.get((key, False, False))
    )
    if filename is None:
        _missing_fonts.add(name or DEFAULT_FONT)
        return None
    path = FONT_DIR / filename
    if not path.is_file():
        _missing_fonts.add(f"{name} ({filename})")
        return None
    cache_key = (str(path), int(round(size_pt * SCALE)))
    if cache_key not in _font_cache:
        _font_cache[cache_key] = ImageFont.truetype(str(path), cache_key[1])
    return _font_cache[cache_key]


def run_style(run, inherited_size: float) -> tuple[str, bool, bool, float]:
    size = run.font.size.pt if run.font.size else inherited_size
    return (run.font.name or DEFAULT_FONT, bool(run.font.bold), bool(run.font.italic), size)


def measure(text: str, style: tuple[str, bool, bool, float]) -> float | None:
    if not text:
        return 0.0
    font = load_font(*style)
    if font is None:
        return None
    return font.getlength(text) / SCALE


def wrapped_line_count(paragraph, usable_pt: float, inherited_size: float) -> tuple[int, float, bool]:
    """Greedy-wrap a paragraph. Returns (lines, widest line, measurable)."""
    pieces: list[tuple[str, tuple[str, bool, bool, float]]] = []
    for run in paragraph.runs:
        style = run_style(run, inherited_size)
        for token in re.split(r"(\s+)", run.text):
            if token:
                pieces.append((token, style))
    if not pieces:
        return 0, 0.0, True

    limit = usable_pt * FIT_TOLERANCE
    lines, widest = 1, 0.0
    current = 0.0
    for token, style in pieces:
        width = measure(token, style)
        if width is None:
            return 0, 0.0, False
        if token.isspace():
            # Trailing space never forces a break.
            current += width
            continue
        if current + width > limit and current > 0:
            widest = max(widest, current)
            lines += 1
            current = width
        else:
            current += width
    widest = max(widest, current)
    return lines, widest, True


def paragraph_size(paragraph, frame_default: float) -> float:
    for run in paragraph.runs:
        if run.font.size:
            return run.font.size.pt
    if paragraph.font.size:
        return paragraph.font.size.pt
    return frame_default


@dataclass
class Rendered:
    """Where a shape's text actually lands, in inches."""
    shape_name: str
    preview: str
    left: float
    right: float
    top: float
    bottom: float
    lines: int
    box_lines: int
    fill: float          # widest single line as % of usable width, 0 if multi-line


def rendered_text(shape) -> Rendered | None:
    frame = shape.text_frame
    text = frame.text.strip()
    if not text:
        return None
    # A lone glyph in a tiny box is decoration (arrows, bullets), not copy.
    if len(text) <= 2 and not any(ch.isalnum() for ch in text):
        return None
    # normAutofit shrinks text to fit; PowerPoint owns that decision, not us.
    if frame.auto_size is not None and "NONE" not in str(frame.auto_size):
        return None
    # A rotated label's real extent is not its stored box, so the geometry below
    # would be wrong. Rotated axis captions read as collisions when they are not.
    if abs(shape.rotation or 0) > 0.01:
        _rotated.add(shape.name)
        return None

    usable_w = (Emu(shape.width).inches * 72
                - (frame.margin_left + frame.margin_right) / 12700.0)
    usable_h = (Emu(shape.height).inches * 72
                - (frame.margin_top + frame.margin_bottom) / 12700.0)
    if usable_w <= 0 or usable_h <= 0:
        return None

    box_left = Emu(shape.left or 0).inches + frame.margin_left / 914400.0
    total_lines = 0
    tallest = 0.0
    best_fill = 0.0
    text_left, text_right = None, None
    for paragraph in frame.paragraphs:
        if not paragraph.runs:
            continue
        size = paragraph_size(paragraph, DEFAULT_SIZE_PT)
        tallest = max(tallest, size * LINE_HEIGHT * (paragraph.line_spacing or 1.0))
        if frame.word_wrap is False:
            widths = [measure(r.text, run_style(r, size)) for r in paragraph.runs]
            if any(w is None for w in widths):
                return None
            lines, widest = 1, sum(widths)
        else:
            lines, widest, ok = wrapped_line_count(paragraph, usable_w, size)
            if not ok:
                return None
        total_lines += lines
        if lines == 1:
            best_fill = max(best_fill, 100 * widest / usable_w)

        # A left-aligned label in a full-width box occupies only its own width.
        # Comparing boxes instead of text makes every side-by-side footer look
        # like a collision, so place the text inside the box by its alignment.
        run_w = min(widest, usable_w) / 72.0
        slack = (usable_w / 72.0) - run_w
        align = str(paragraph.alignment or "LEFT")
        if "CENTER" in align:
            start = box_left + slack / 2
        elif "RIGHT" in align:
            start = box_left + slack
        else:
            start = box_left
        text_left = start if text_left is None else min(text_left, start)
        text_right = (start + run_w) if text_right is None else max(text_right, start + run_w)

    if total_lines == 0 or tallest == 0 or text_left is None:
        return None

    block = total_lines * tallest / 72.0            # inches
    top_in = Emu(shape.top or 0).inches + frame.margin_top / 914400.0
    box_h = usable_h / 72.0
    anchor = str(frame.vertical_anchor or "TOP")
    if "MIDDLE" in anchor:
        start = top_in + (box_h - block) / 2
    elif "BOTTOM" in anchor:
        start = top_in + box_h - block
    else:
        start = top_in

    return Rendered(
        shape_name=shape.name,
        preview=text.replace("\n", " ")[:58],
        left=text_left,
        right=text_right,
        top=start,
        bottom=start + block,
        lines=total_lines,
        box_lines=max(1, int(math.floor(usable_h / tallest + 0.01))),
        fill=best_fill if total_lines == 1 else 0.0,
    )


def check_deck(pptx: Path, warn_fill: float) -> list[Finding]:
    prs = Presentation(str(pptx))
    total = len(prs.slides)
    canvas = (Emu(prs.slide_width).inches, Emu(prs.slide_height).inches)
    findings: list[Finding] = []
    numbered: list[int] = []

    for index, slide in enumerate(prs.slides, 1):
        blocks: list[Rendered] = []
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            block = rendered_text(shape)
            if block is not None:
                blocks.append(block)

            text = shape.text_frame.text.strip()
            match = PAGE_NUMBER.search(text)
            if match:
                numbered.append(index)
                shown, out_of = int(match.group(1)), int(match.group(2))
                if shown != index or out_of != total:
                    findings.append(Finding(
                        index, "numbering",
                        "footer reads %d / %d, expected %d / %d" % (shown, out_of, index, total),
                    ))

        # Rendered text landing on top of other rendered text.
        for i, a in enumerate(blocks):
            for b in blocks[i + 1:]:
                overlap_y = min(a.bottom, b.bottom) - max(a.top, b.top)
                overlap_x = min(a.right, b.right) - max(a.left, b.left)
                if overlap_y > 0.01 and overlap_x > 0.05:
                    findings.append(Finding(
                        index, "collision",
                        "%r (%d line%s) overlaps %r by %.2fin vertically%s"
                        % (a.preview, a.lines, "" if a.lines == 1 else "s",
                           b.preview, overlap_y,
                           "" if overlap_y >= COLLISION_FATAL_IN else " -- near the"
                           " model's margin of error, render the slide to confirm"),
                        fatal=overlap_y >= COLLISION_FATAL_IN,
                    ))

        for block in blocks:
            if block.top < -0.01 or block.bottom > canvas[1] + 0.01:
                findings.append(Finding(
                    index, "off-canvas",
                    "text spans %.2f-%.2fin on a %.2fin-tall canvas: %r"
                    % (block.top, block.bottom, canvas[1], block.preview),
                ))
            elif block.lines > block.box_lines:
                findings.append(Finding(
                    index, "spills",
                    "%d lines in a box sized for %d: %r"
                    % (block.lines, block.box_lines, block.preview),
                    fatal=False,
                ))
            if block.fill >= warn_fill:
                findings.append(Finding(
                    index, "tight",
                    "single line at %.0f%% of box width: %r" % (block.fill, block.preview),
                    fatal=False,
                ))

        bodies = [(sh.text_frame.text, "text") for sh in slide.shapes if sh.has_text_frame]
        if slide.has_notes_slide:
            bodies.append((slide.notes_slide.notes_text_frame.text, "notes"))
        for body, where in bodies:
            for ref in SLIDE_REF.finditer(body):
                target = int(ref.group(1))
                if not 1 <= target <= total:
                    findings.append(Finding(
                        index, "cross-ref",
                        "%s points at slide %d; deck has %d" % (where, target, total),
                    ))
                else:
                    findings.append(Finding(
                        index, "cross-ref",
                        "%s points at slide %d -- confirm the target after any reorder"
                        % (where, target),
                        fatal=False,
                    ))

    missing = [n for n in range(1, total + 1) if n not in numbered]
    if numbered and missing:
        findings.append(Finding(
            0, "numbering",
            "%d of %d slides carry an 'N / M' footer; missing on %s"
            % (len(numbered), total,
               ", ".join(str(n) for n in missing[:12]) + (" ..." if len(missing) > 12 else "")),
            fatal=False,
        ))
    return findings


def report(pptx: Path, findings: list[Finding]) -> bool:
    fatal = [f for f in findings if f.fatal]
    warn = [f for f in findings if not f.fatal]
    print("=" * 78)
    print(pptx.name)
    if not findings:
        print("  clean")
    for group, label in ((fatal, "FAIL"), (warn, "note")):
        for f in group:
            where = "slide %2d" % f.slide if f.slide else "deck    "
            print("  %-4s %s  %-10s %s" % (label, where, f.kind, f.detail))
    return not fatal


def main(argv: list[str] | None = None) -> int:
    configure_utf8_stdio()
    parser = argparse.ArgumentParser(
        description="Check companion deck slides for overflow, stale numbering, "
                    "and stale cross-references.",
    )
    parser.add_argument("pptx", nargs="?", type=Path, help="Path to the .pptx file")
    parser.add_argument("--study", help="Study slug under Studies/")
    parser.add_argument("--deck", help="Deck filename under Studies/<Slug>/")
    parser.add_argument("--all", action="store_true", help="Check every deck under Studies/")
    parser.add_argument(
        "--warn-fill", type=float, default=95.0,
        help="Flag single-line boxes at or above this %% of box width (default: 95)",
    )
    args = parser.parse_args(argv)

    if args.all:
        decks = sorted(STUDIES.glob("*/*.pptx"))
        if not decks:
            raise SystemExit("No decks found under Studies/")
    else:
        decks = [resolve_pptx(args.pptx, args.study, args.deck)]

    ok = True
    for deck in decks:
        ok &= report(deck, check_deck(deck, args.warn_fill))
    if _missing_fonts:
        print("\nNot measured (no font file): %s" % ", ".join(sorted(_missing_fonts)))
    if _rotated:
        print("\nNot measured (%d rotated shape%s): geometry assumes an unrotated box"
              % (len(_rotated), "" if len(_rotated) == 1 else "s"))
    return 0 if ok else 1


if __name__ == "__main__":
    raise SystemExit(main())

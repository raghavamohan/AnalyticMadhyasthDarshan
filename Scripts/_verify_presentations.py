#!/usr/bin/env python3
"""Verify the manifest and generated artifacts for companion presentations."""
from __future__ import annotations

import argparse
import re
from collections import Counter
from pathlib import Path

import pymupdf as fitz
from pptx import Presentation

from _common import BASE, configure_utf8_stdio
from _presentation_pipeline import (
    DeckSpec,
    PresentationManifest,
    load_manifest,
    manifest_errors,
    repo_relative,
)

TOKEN_RE = re.compile(r"[^\W_]+", re.UNICODE)
SUBSET_PREFIX_RE = re.compile(r"^[A-Z]{6}\+")
MIN_TEXT_RECALL = 0.98
ASPECT_RATIO_TOLERANCE = 0.002


def tokens(text: str) -> Counter[str]:
    return Counter(TOKEN_RE.findall(text.casefold()))


def token_recall(expected: str, actual: str) -> float:
    wanted = tokens(expected)
    if not wanted:
        return 1.0
    found = tokens(actual)
    matched = sum(min(count, found[token]) for token, count in wanted.items())
    return matched / sum(wanted.values())


def slide_text(slide) -> str:
    return "\n".join(
        shape.text_frame.text
        for shape in slide.shapes
        if getattr(shape, "has_text_frame", False) and shape.text_frame.text.strip()
    )


def note_text(slide) -> str:
    if not slide.has_notes_slide:
        return ""
    return slide.notes_slide.notes_text_frame.text.strip()


def _font_family(pdf_name: str) -> str:
    name = SUBSET_PREFIX_RE.sub("", pdf_name)
    name = re.sub(r"[ _-](bold|italic|oblique|regular|roman).*$", "", name, flags=re.I)
    return re.sub(r"\s+", "", name).casefold()


def pdf_font_families(doc: fitz.Document) -> set[str]:
    families: set[str] = set()
    for page in doc:
        for font in page.get_fonts(full=True):
            families.add(_font_family(font[3]))
    return families


def page_has_content(page: fitz.Page) -> bool:
    return bool(page.get_text("text").strip() or page.get_images() or page.get_drawings())


def _generated_path(configured: Path, output_root: Path | None) -> Path:
    if output_root is None:
        return configured
    return output_root.resolve() / configured.resolve().relative_to(BASE.resolve())


def verify_slides_pdf(spec: DeckSpec, prs: Presentation, pdf_path: Path) -> list[str]:
    errors: list[str] = []
    if not pdf_path.is_file():
        return [f"{spec.id}: slides PDF is missing: {pdf_path}"]
    try:
        doc = fitz.open(pdf_path)
    except Exception as exc:  # noqa: BLE001 - verifier boundary
        return [f"{spec.id}: cannot open slides PDF: {exc}"]

    try:
        total = len(prs.slides)
        if len(doc) != total:
            errors.append(f"{spec.id}: slides PDF has {len(doc)} pages; PPTX has {total} slides")

        expected_ratio = float(prs.slide_width) / float(prs.slide_height)
        for index, page in enumerate(doc, 1):
            actual_ratio = page.rect.width / page.rect.height
            if abs(actual_ratio - expected_ratio) > ASPECT_RATIO_TOLERANCE:
                errors.append(
                    f"{spec.id}: slides PDF page {index} aspect ratio {actual_ratio:.4f}; "
                    f"PPTX is {expected_ratio:.4f}"
                )
            if not page_has_content(page):
                errors.append(f"{spec.id}: slides PDF page {index} is blank")

        for index, (slide, page) in enumerate(zip(prs.slides, doc), 1):
            recall = token_recall(slide_text(slide), page.get_text("text"))
            if recall < MIN_TEXT_RECALL:
                errors.append(
                    f"{spec.id}: slides PDF page {index} preserves only {recall:.1%} "
                    "of PPTX text"
                )

        families = pdf_font_families(doc)
        for required in spec.required_fonts:
            if _font_family(required) not in families:
                errors.append(
                    f"{spec.id}: slides PDF does not embed required font family {required}; "
                    f"found {', '.join(sorted(families)) or 'none'}"
                )
    finally:
        doc.close()
    return errors


def verify_notes_pdf(spec: DeckSpec, prs: Presentation, pdf_path: Path) -> list[str]:
    errors: list[str] = []
    if not pdf_path.is_file():
        return [f"{spec.id}: notes PDF is missing: {pdf_path}"]
    try:
        doc = fitz.open(pdf_path)
    except Exception as exc:  # noqa: BLE001 - verifier boundary
        return [f"{spec.id}: cannot open notes PDF: {exc}"]

    try:
        total = len(prs.slides)
        if len(doc) < total:
            errors.append(f"{spec.id}: notes PDF has {len(doc)} pages; PPTX has {total} slides")
        all_text = "\n".join(page.get_text("text") for page in doc)
        expected_notes = "\n".join(note_text(slide) for slide in prs.slides)
        recall = token_recall(expected_notes, all_text)
        if recall < MIN_TEXT_RECALL:
            errors.append(
                f"{spec.id}: notes PDF preserves only {recall:.1%} of speaker-note text"
            )
        folded = re.sub(r"\s+", " ", all_text).casefold()
        for index in range(1, total + 1):
            marker = f"slide {index} of {total}"
            if marker not in folded:
                errors.append(f"{spec.id}: notes PDF is missing header {marker!r}")
        blank = [str(index) for index, page in enumerate(doc, 1) if not page_has_content(page)]
        if blank:
            errors.append(f"{spec.id}: notes PDF has blank pages: {', '.join(blank)}")

        families = pdf_font_families(doc)
        for required in spec.required_fonts:
            if _font_family(required) not in families:
                errors.append(
                    f"{spec.id}: notes PDF does not embed required font family {required}; "
                    f"found {', '.join(sorted(families)) or 'none'}"
                )
    finally:
        doc.close()
    return errors


def verify_deck(spec: DeckSpec, output_root: Path | None = None) -> list[str]:
    if not spec.source.is_file():
        return [f"{spec.id}: PPTX is missing: {repo_relative(spec.source)}"]
    prs = Presentation(str(spec.source))
    errors: list[str] = []
    if not prs.slides:
        errors.append(f"{spec.id}: PPTX has no slides")
    missing_notes = [
        str(index)
        for index, slide in enumerate(prs.slides, 1)
        if not note_text(slide)
    ]
    if missing_notes:
        errors.append(f"{spec.id}: slides without speaker notes: {', '.join(missing_notes)}")
    errors.extend(verify_slides_pdf(spec, prs, _generated_path(spec.slides_pdf, output_root)))
    errors.extend(verify_notes_pdf(spec, prs, _generated_path(spec.notes_pdf, output_root)))
    return errors


def verify_manifest_only(manifest: PresentationManifest) -> list[str]:
    return manifest_errors(manifest)


def main(argv: list[str] | None = None) -> int:
    configure_utf8_stdio()
    parser = argparse.ArgumentParser(
        description="Verify presentation manifest coverage and generated slides/notes PDFs."
    )
    parser.add_argument(
        "--manifest", type=Path, default=None,
        help="Manifest path (default: Scripts/presentation-pipeline.json)",
    )
    parser.add_argument(
        "--manifest-only", action="store_true",
        help="Check coverage/path safety without requiring generated PDFs",
    )
    parser.add_argument(
        "--deck", action="append", default=[], metavar="ID",
        help="Verify one manifest deck id; repeat for multiple (default: all)",
    )
    parser.add_argument(
        "--output-root", type=Path,
        help="Root containing generated repository-relative PDF paths",
    )
    args = parser.parse_args(argv)

    try:
        manifest = load_manifest(args.manifest) if args.manifest else load_manifest()
    except (OSError, ValueError, TypeError) as exc:
        print(f"Presentation manifest invalid: {exc}")
        return 1

    errors = verify_manifest_only(manifest)
    if not args.manifest_only:
        selected = manifest.decks
        if args.deck:
            try:
                selected = tuple(manifest.deck(deck_id) for deck_id in args.deck)
            except ValueError as exc:
                errors.append(str(exc))
                selected = ()
        for spec in selected:
            errors.extend(verify_deck(spec, args.output_root))

    if errors:
        print("Presentation verification FAILED:")
        for error in errors:
            print(f"  - {error}")
        return 1
    scope = "manifest" if args.manifest_only else (
        ", ".join(args.deck) if args.deck else "all presentations"
    )
    print(f"Presentation verification passed ({scope}).")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())

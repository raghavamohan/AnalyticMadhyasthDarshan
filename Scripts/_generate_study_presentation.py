#!/usr/bin/env python3
"""Generate a study presentation (.pptx) from curated slides YAML.

Default paths for The-Ontology-of-Coexistence:
  Studies/<Slug>/<Slug>-ontology-slides.yaml → <Slug>-ontology.pptx

Requires: python-pptx, PyYAML, npm install in Scripts/ (Puppeteer for SVG figures).
"""
from __future__ import annotations

import argparse
import subprocess
import sys
import tempfile
from pathlib import Path

import yaml
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

REPO = Path(__file__).resolve().parent.parent
SCRIPTS = REPO / "Scripts"
STUDIES = REPO / "Studies"
FOOTER_TEXT = "AnalyticMadhyasthDarshan.org"


def default_slides_yaml(slug: str) -> Path:
    return STUDIES / slug / f"{slug}-ontology-slides.yaml"


def default_output_pptx(slug: str) -> Path:
    return STUDIES / slug / f"{slug}-ontology.pptx"


def rasterize_svg(svg_path: Path, cache_dir: Path) -> Path:
    png_path = cache_dir / f"{svg_path.stem}.png"
    if png_path.exists() and png_path.stat().st_mtime >= svg_path.stat().st_mtime:
        return png_path
    subprocess.run(
        ["node", str(SCRIPTS / "_svg_to_png.js"), str(svg_path), str(png_path)],
        check=True,
        cwd=REPO,
    )
    return png_path


def _set_run_font(run, size_pt: int = 18, bold: bool = False) -> None:
    run.font.size = Pt(size_pt)
    run.font.bold = bold
    run.font.color.rgb = RGBColor(0x1A, 0x1A, 0x1A)


def _add_textbox(
    slide,
    left,
    top,
    width,
    height,
    text: str,
    size_pt: int = 18,
    bold: bool = False,
    align=PP_ALIGN.LEFT,
) -> None:
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    _set_run_font(run, size_pt, bold)


def _add_bullets(slide, left, top, width, height, bullets: list[str], size_pt: int = 18) -> None:
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    for i, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.level = 0
        p.space_after = Pt(6)
        run = p.add_run()
        run.text = bullet
        _set_run_font(run, size_pt)


def _add_quote_block(
    slide, left, top, width, height, quote: str, citation: str | None, size_pt: int = 14
) -> None:
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = quote
    _set_run_font(run, size_pt)
    run.font.italic = True
    if citation:
        p2 = tf.add_paragraph()
        p2.space_before = Pt(4)
        run2 = p2.add_run()
        run2.text = citation
        _set_run_font(run2, 12)
        run2.font.color.rgb = RGBColor(0x55, 0x55, 0x55)


def _add_footer(slide, slide_num: int, total: int) -> None:
    footer_y = Inches(7.05)
    _add_textbox(
        slide,
        Inches(0.5),
        footer_y,
        Inches(8),
        Inches(0.35),
        FOOTER_TEXT,
        size_pt=10,
    )
    _add_textbox(
        slide,
        Inches(11.5),
        footer_y,
        Inches(1.5),
        Inches(0.35),
        f"{slide_num} / {total}",
        size_pt=10,
        align=PP_ALIGN.RIGHT,
    )


def _add_speaker_notes(slide, notes: str | None) -> None:
    if not notes:
        return
    notes_slide = slide.notes_slide
    tf = notes_slide.notes_text_frame
    tf.text = notes


def _add_figure(slide, png_path: Path, left, top, max_width, max_height) -> None:
    pic = slide.shapes.add_picture(str(png_path), left, top)
    scale_w = max_width / pic.width if pic.width > max_width else 1.0
    scale_h = max_height / pic.height if pic.height > max_height else 1.0
    scale = min(scale_w, scale_h)
    if scale < 1.0:
        pic.width = int(pic.width * scale)
        pic.height = int(pic.height * scale)


def build_presentation(slides_data: dict, study_dir: Path, cache_dir: Path) -> Presentation:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    slide_defs = slides_data.get("slides", [])
    total = len(slide_defs)

    for index, slide_def in enumerate(slide_defs, start=1):
        layout_type = slide_def.get("layout", "content")
        title = slide_def.get("title", "")
        bullets = slide_def.get("bullets") or []
        quote = slide_def.get("quote")
        citation = slide_def.get("citation")
        figure = slide_def.get("figure")
        notes = slide_def.get("speaker_notes")

        blank = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank)

        margin = Inches(0.55)
        content_top = Inches(1.15)
        content_bottom = Inches(6.85)
        full_width = prs.slide_width - margin * 2

        if layout_type == "title":
            _add_textbox(
                slide,
                margin,
                Inches(2.0),
                full_width,
                Inches(1.2),
                title,
                size_pt=36,
                bold=True,
                align=PP_ALIGN.CENTER,
            )
            subtitle = slide_def.get("subtitle", "")
            if subtitle:
                _add_textbox(
                    slide,
                    margin,
                    Inches(3.3),
                    full_width,
                    Inches(1.5),
                    subtitle,
                    size_pt=22,
                    align=PP_ALIGN.CENTER,
                )
        else:
            _add_textbox(
                slide,
                margin,
                Inches(0.35),
                full_width,
                Inches(0.75),
                title,
                size_pt=28,
                bold=True,
            )

            figure_png: Path | None = None
            if figure:
                svg_path = study_dir / figure
                if not svg_path.is_file():
                    raise FileNotFoundError(f"Figure not found: {svg_path}")
                figure_png = rasterize_svg(svg_path, cache_dir)

            has_figure = figure_png is not None
            text_width = Inches(6.2) if has_figure else full_width
            text_height = content_bottom - content_top

            if bullets:
                _add_bullets(slide, margin, content_top, text_width, text_height, bullets)

            quote_top = content_top
            if bullets:
                quote_top = content_top + Inches(min(0.35 * len(bullets) + 0.2, 3.5))
            if quote:
                _add_quote_block(
                    slide,
                    margin,
                    quote_top,
                    text_width,
                    Inches(2.0),
                    quote,
                    citation,
                )

            if figure_png:
                fig_left = Inches(7.0)
                fig_top = content_top
                fig_max_w = Inches(5.9)
                fig_max_h = Inches(5.5)
                _add_figure(slide, figure_png, fig_left, fig_top, fig_max_w, fig_max_h)

        _add_footer(slide, index, total)
        _add_speaker_notes(slide, notes)

    return prs


def main(argv: list[str] | None = None) -> int:
    parser = argparse.ArgumentParser(description="Generate study presentation from YAML.")
    parser.add_argument("slug", help="Study directory name under Studies/")
    parser.add_argument(
        "--slides",
        type=Path,
        help="Path to slides YAML (default: Studies/<Slug>/<Slug>-ontology-slides.yaml)",
    )
    parser.add_argument(
        "--output",
        type=Path,
        help="Output .pptx path (default: Studies/<Slug>/<Slug>-ontology.pptx)",
    )
    args = parser.parse_args(argv)

    slug = args.slug
    study_dir = STUDIES / slug
    if not study_dir.is_dir():
        print(f"Study directory not found: {study_dir}", file=sys.stderr)
        return 1

    slides_yaml = args.slides or default_slides_yaml(slug)
    output = args.output or default_output_pptx(slug)

    if not slides_yaml.is_file():
        print(f"Slides YAML not found: {slides_yaml}", file=sys.stderr)
        return 1

    with slides_yaml.open(encoding="utf-8") as handle:
        slides_data = yaml.safe_load(handle)

    with tempfile.TemporaryDirectory() as tmp:
        cache_dir = Path(tmp)
        prs = build_presentation(slides_data, study_dir, cache_dir)
        output.parent.mkdir(parents=True, exist_ok=True)
        prs.save(str(output))

    print(f"Wrote {output}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())

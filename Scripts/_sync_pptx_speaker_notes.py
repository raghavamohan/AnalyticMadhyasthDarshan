#!/usr/bin/env python3
"""Sync PowerPoint speaker notes from a JSON map of slide number → note text.

The JSON object keys are slide numbers as strings or integers (1-based).
Values are plain speaker-note strings for the notes pane.

Examples (from repo root):

  python Scripts/_sync_pptx_speaker_notes.py Studies/.../Deck.pptx Studies/.../Deck.notes.json
  python Scripts/_sync_pptx_speaker_notes.py --pptx Studies/.../Deck.pptx --notes Studies/.../notes.json
"""
from __future__ import annotations

import argparse
import json
import sys
from pathlib import Path

from pptx import Presentation


def load_notes(path: Path) -> dict[int, str]:
    data = json.loads(path.read_text(encoding="utf-8"))
    if not isinstance(data, dict):
        raise SystemExit(f"Notes JSON must be an object: {path}")
    notes: dict[int, str] = {}
    for key, value in data.items():
        try:
            idx = int(key)
        except (TypeError, ValueError) as exc:
            raise SystemExit(f"Invalid slide key {key!r} in {path}") from exc
        if idx < 1:
            raise SystemExit(f"Slide numbers must be >= 1; got {idx}")
        if not isinstance(value, str):
            raise SystemExit(f"Notes for slide {idx} must be a string")
        notes[idx] = value.strip()
    if not notes:
        raise SystemExit(f"No notes found in {path}")
    return notes


def sync_speaker_notes(pptx: Path, notes: dict[int, str]) -> None:
    pptx = pptx.expanduser().resolve()
    if not pptx.is_file():
        raise SystemExit(f"PPTX not found: {pptx}")

    prs = Presentation(str(pptx))
    slide_count = len(prs.slides)
    missing = [i for i in range(1, slide_count + 1) if i not in notes]
    extras = sorted(i for i in notes if i > slide_count)
    if missing:
        raise SystemExit(
            f"{pptx.name} has {slide_count} slides but notes JSON is missing: "
            + ", ".join(str(i) for i in missing)
        )
    if extras:
        raise SystemExit(
            f"Notes JSON has slides beyond deck length ({slide_count}): "
            + ", ".join(str(i) for i in extras)
        )

    for idx, slide in enumerate(prs.slides, 1):
        notes_slide = slide.notes_slide
        tf = notes_slide.notes_text_frame
        tf.clear()
        # Split on blank lines into paragraphs so PowerPoint does not turn
        # newlines into vertical-tab characters in a single paragraph.
        paragraphs = [
            " ".join(block.split())
            for block in notes[idx].replace("\r\n", "\n").split("\n\n")
            if block.strip()
        ]
        if not paragraphs:
            paragraphs = [""]
        tf.paragraphs[0].text = paragraphs[0]
        for para_text in paragraphs[1:]:
            p = tf.add_paragraph()
            p.text = para_text
    prs.save(str(pptx))


def main(argv: list[str] | None = None) -> int:
    parser = argparse.ArgumentParser(
        description="Write speaker notes into a .pptx from a JSON slide map."
    )
    parser.add_argument(
        "pptx_pos",
        nargs="?",
        type=Path,
        help="Path to the .pptx deck",
    )
    parser.add_argument(
        "notes_pos",
        nargs="?",
        type=Path,
        help="Path to the notes JSON map",
    )
    parser.add_argument("--pptx", type=Path, help="Path to the .pptx deck")
    parser.add_argument("--notes", type=Path, help="Path to the notes JSON map")
    args = parser.parse_args(argv)

    pptx = args.pptx or args.pptx_pos
    notes_path = args.notes or args.notes_pos
    if pptx is None or notes_path is None:
        parser.error("Provide both a .pptx path and a notes JSON path")

    notes = load_notes(notes_path.expanduser().resolve())
    sync_speaker_notes(pptx, notes)
    print(
        f"Updated speaker notes for {len(notes)} slides in {pptx.expanduser().resolve()}"
    )
    return 0


if __name__ == "__main__":
    sys.exit(main())
